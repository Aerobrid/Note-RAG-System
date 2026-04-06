"""
Parser — uses IBM Docling to convert PDF, DOCX, PPTX to markdown,
and handles code files as structured text with language detection.
"""
from __future__ import annotations

import re
from pathlib import Path
from dataclasses import dataclass, field

import pdfplumber

CODE_EXTENSIONS = {
    ".py": "python", ".js": "javascript", ".ts": "typescript",
    ".jsx": "jsx", ".tsx": "tsx", ".java": "java", ".c": "c",
    ".cpp": "c++", ".h": "c", ".cs": "csharp", ".go": "go",
    ".rs": "rust", ".rb": "ruby", ".php": "php", ".swift": "swift",
    ".kt": "kotlin", ".sql": "sql", ".sh": "bash", ".r": "r",
    ".m": "matlab", ".ipynb": "notebook",
}

DOC_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md", ".html"}


@dataclass
class ParsedDocument:
    text: str
    source: str
    doc_type: str          # "document" | "code"
    language: str = ""     # for code files
    title: str = ""
    page_count: int = 0
    metadata: dict = field(default_factory=dict)


def parse_file(file_path: str | Path) -> ParsedDocument:
    """Parse any supported file type and return structured text."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix in CODE_EXTENSIONS:
        return _parse_code(path)
    elif suffix == ".ipynb":
        return _parse_notebook(path)
    elif suffix in DOC_EXTENSIONS:
        return _parse_document(path)
    else:
        # Fallback: treat as plain text
        try:
            text = path.read_text(encoding="utf-8", errors="replace")
            return ParsedDocument(
                text=text, source=str(path),
                doc_type="document", title=path.stem,
            )
        except Exception as e:
            raise ValueError(f"Cannot parse {path.name}: {e}") from e


def _parse_document(path: Path) -> ParsedDocument:
    """Parse PDF using pdfplumber, fallback to plain text for others."""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        with pdfplumber.open(str(path)) as pdf:
            text = "\n".join(page.extract_text() or "" for page in pdf.pages)
        page_count = len(pdf.pages)
        return ParsedDocument(
            text=text,
            source=str(path),
            doc_type="document",
            title=path.stem.replace("_", " ").replace("-", " ").title(),
            page_count=page_count,
            metadata={
                "source": str(path),
                "filename": path.name,
                "file_type": path.suffix.lower(),
                "title": path.stem,
            },
        )
    text = ""
    page_count = 0
    # PowerPoint (.pptx) support (lazy import)
    if suffix == ".pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            parts: list[str] = []
            for slide in prs.slides:
                slide_text: list[str] = []
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text.strip())
                    except Exception:
                        continue
                if slide_text:
                    parts.append("\n".join(slide_text))

            text = "\n\n".join(parts).strip()
            page_count = len(prs.slides)
            if not text:
                raise ValueError("No text extracted from pptx")
        except Exception:
            pass

    # DOCX support (lazy import)
    elif suffix == ".docx":
        try:
            import docx

            doc = docx.Document(str(path))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
            text = "\n\n".join(paragraphs)
            page_count = len(paragraphs)
            if not text:
                raise ValueError("No text extracted from docx")
        except Exception:
            pass

    if text:
        return ParsedDocument(
            text=text,
            source=str(path),
            doc_type="document",
            title=path.stem.replace("_", " ").replace("-", " ").title() if suffix == ".pdf" else path.stem,
            page_count=page_count,
            metadata={
                "source": str(path),
                "filename": path.name,
                "file_type": path.suffix.lower(),
                "title": path.stem,
            },
        )

    # Fallback: treat as plain text for other types (including .ppt legacy)
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return ParsedDocument(
            text=text, source=str(path),
            doc_type="document", title=path.stem,
        )
    except Exception as e:
        raise ValueError(f"Cannot parse {path.name}: {e}") from e


def _parse_code(path: Path) -> ParsedDocument:
    """Parse a code file, adding language metadata and structure comments."""
    language = CODE_EXTENSIONS.get(path.suffix.lower(), "text")
    raw = path.read_text(encoding="utf-8", errors="replace")

    # Add a structured header so the LLM understands what it's reading
    structured = (
        f"# File: {path.name}\n"
        f"# Language: {language}\n"
        f"# Path: {path}\n\n"
        f"```{language}\n{raw}\n```"
    )

    return ParsedDocument(
        text=structured,
        source=str(path),
        doc_type="code",
        language=language,
        title=path.name,
        metadata={
            "source": str(path),
            "filename": path.name,
            "file_type": path.suffix.lower(),
            "language": language,
            "title": path.name,
            "doc_type": "code",
        },
    )


def _parse_notebook(path: Path) -> ParsedDocument:
    """Parse Jupyter notebooks — extract code cells and markdown cells."""
    import json

    with path.open(encoding="utf-8") as f:
        nb = json.load(f)

    parts: list[str] = [f"# Notebook: {path.name}\n"]
    for i, cell in enumerate(nb.get("cells", []), 1):
        cell_type = cell.get("cell_type", "")
        source = "".join(cell.get("source", []))
        if not source.strip():
            continue
        if cell_type == "markdown":
            parts.append(source)
        elif cell_type == "code":
            parts.append(f"```python\n{source}\n```")

    text = "\n\n".join(parts)
    return ParsedDocument(
        text=text,
        source=str(path),
        doc_type="code",
        language="python",
        title=path.stem,
        metadata={
            "source": str(path),
            "filename": path.name,
            "file_type": ".ipynb",
            "language": "python",
            "doc_type": "code",
        },
    )
